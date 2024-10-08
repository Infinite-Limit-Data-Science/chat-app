from typing import Union, Iterator, List
from pathlib import Path
from pptx import Presentation
from pptx.shapes.shapetree import SlideShapes
from pptx.text.text import TextFrame
from pptx.table import Table
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_core.documents import Document
from orchestrators.doc.document_loaders.base_loader import BaseLoader

class PowerPointLoader(BaseLoader):
    """
    Recursively extract all text from all slides in presentation, including:
    - textframes
    - tables
    - pictures
    - group shapes

    Note tables, pictures, and group shapes do not directly have textframes.
    Recursive scans on shapetree are handled appropriately to capture text
    """
    def __init__(self, file_path: Union[str, Path]):
        self._file_path = file_path
        self.doc = Presentation(self._file_path)

    def lazy_load(self) -> Iterator[Document]:
        """Lazily load document"""
        full_text = []

        for slide in self.doc.slides:
            full_text.extend(self._extract_text_from_shapes(slide.shapes))
        
        metadata = {'source': str(self._file_path)}
        yield Document(page_content="\n".join(full_text), metadata=metadata)
    load = lazy_load

    def _extract_text_from_shapes(self, shapes: SlideShapes) -> List[str]:
        """Extract text from all shapes in a slide"""
        text_parts = []
        for shape in shapes:
            if shape.has_text_frame:
                text_parts.append(self._extract_text_from_text_frame(shape))
            if shape.has_table:
                text_parts.append(self._extract_text_from_table(shape.table))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                text_parts.append(self._extract_text_from_image(shape)) 
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                text_parts.extend(self._extract_text_from_shapes(shape.shapes))
        return text_parts

    def _extract_text_from_text_frame(self, text_frame: TextFrame) -> str:
        """Extract text from text frames (titles, text boxes, etc.)"""
        return text_frame.text

    def _extract_text_from_table(self, table: Table) -> str:
        """Extract text from table cells (in PPT tables cannot have nested tables)"""
        table_text = []
        for row in table.rows:
            for cell in row.cells:
                table_text.append(cell.text)
        return "\n".join(table_text)

    def _extract_text_from_image(self, _) -> str:
        return ''